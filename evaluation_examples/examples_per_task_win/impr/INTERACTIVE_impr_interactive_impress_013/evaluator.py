import logging
import os
import re

logger = logging.getLogger("desktopenv.metric.general")


def check_pptx_author_info(result_path, **options):
    """
    Check that the PPTX/ODP file contains author-related information.

    The task instruction is "Help me add author info to the slides" — it does
    not prescribe a specific author name.  This evaluator therefore looks for
    *any* plausible author / contact information that has been added to the
    presentation, rather than checking for a hard-coded name.

    Args:
        result_path: Local filesystem path to the presentation file (pulled
                     from the VM via the ``vm_file`` result getter).
        **options:  Reserved for framework compatibility; currently unused.

    Returns:
        float: 1.0 if author-related content is detected, 0.0 otherwise.
    """
    if result_path is None or not os.path.exists(result_path):
        logger.warning("Result file not found: %s", result_path)
        return 0.0

    text = _extract_presentation_text(result_path)
    if not text or not text.strip():
        logger.warning("No text extracted from presentation file: %s", result_path)
        return 0.0

    logger.info("Extracted %d characters of text from presentation", len(text))

    # ------------------------------------------------------------------
    # Patterns that indicate author / presenter / contact information.
    # The list is deliberately broad so that any reasonable author info
    # added by the agent is accepted.
    # ------------------------------------------------------------------
    author_patterns = [
        # Title + name:  Dr. Smith, Professor Jones, Prof. Lee
        r'Dr\.\s+\w+',
        r'Professor\s+\w+',
        r'Prof\.\s+\w+',
        # Email address (standard pattern)
        r'[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}',
        # Institutional affiliation clues
        r'Department\s+of\s+\w+',
        r'University\s+of\s+\w+',
        r'Institute\s+of\s+\w+',
        # Explicit author / presenter labels
        r'Author\s*:',
        r'Prepared\s+by\s*:',
        r'Presented\s+by\s*:',
        r'By\s*:',
    ]

    for pattern in author_patterns:
        if re.search(pattern, text, re.IGNORECASE):
            logger.info("Author-info pattern matched: %s", pattern)
            return 1.0

    # Broader fallback: any plausible full-name pattern (two consecutive
    # capitalized words) that might represent a person's name added as
    # author metadata.
    name_pattern = r'\b[A-Z][a-z]{2,}\s+[A-Z][a-z]{2,}\b'
    if re.search(name_pattern, text):
        logger.info("Fallback name-like pattern matched in presentation text")
        return 1.0

    logger.warning("No author-info patterns found in presentation text")
    return 0.0


def _extract_presentation_text(filepath):
    """
    Extract all visible text from a PPTX (OOXML) presentation.

    Tries python-pptx first; falls back to raw ZIP + XML parsing if
    python-pptx is unavailable or raises an exception.
    """
    # --- primary path: python-pptx ------------------------------------
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text:
                                text_parts.append(run.text)
        result = ' '.join(text_parts)
        if result.strip():
            return result
    except Exception:
        logger.debug("python-pptx extraction failed, trying XML fallback", exc_info=True)

    # --- fallback: ZIP + raw XML (PPTX is a ZIP archive) --------------
    try:
        import zipfile
        from xml.etree import ElementTree

        text_parts = []
        with zipfile.ZipFile(filepath) as zf:
            for name in zf.namelist():
                # PPTX slide files: ppt/slides/slideN.xml
                if name.startswith('ppt/slides/slide') and name.endswith('.xml'):
                    xml_bytes = zf.read(name)
                    root = ElementTree.fromstring(xml_bytes)
                    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    for t_elem in root.iter(f'{{{ns}}}t'):
                        if t_elem.text:
                            text_parts.append(t_elem.text)
        return ' '.join(text_parts)
    except Exception:
        logger.debug("XML fallback extraction also failed", exc_info=True)

    return ""
