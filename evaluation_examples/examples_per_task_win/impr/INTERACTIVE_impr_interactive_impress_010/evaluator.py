import logging
from typing import Dict, List

logger = logging.getLogger("desktopenv.metric.general")


def check_pptx_text(result_path: str, rules: Dict[str, List[str]], **options) -> float:
    """Read all text from a .pptx file and verify include/exclude rules.

    Pulls text from all shapes with text frames across all slides, then checks
    that every string in rules['include'] appears in the extracted text and no
    string in rules['exclude'] appears.

    Args:
        result_path: Local filesystem path to the .pptx file.
        rules: Dict with 'include' (list of substrings that must be present)
               and 'exclude' (list of substrings that must be absent).

    Returns:
        1.0 if all include substrings are found and no exclude substrings are
        found; otherwise 0.0.
    """
    if result_path is None:
        logger.warning("result_path is None, returning 0.0")
        return 0.

    try:
        from pptx import Presentation
    except ImportError as e:
        logger.error(f"python-pptx not available: {e}")
        return 0.

    try:
        prs = Presentation(result_path)
        text = ' '.join(
            run.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            for run in para.runs
        )
    except Exception as e:
        logger.error(f"Failed to read PPTX file '{result_path}': {e}")
        return 0.

    include = rules.get("include", [])
    exclude = rules.get("exclude", [])

    if all(r in text for r in include) and all(r not in text for r in exclude):
        return 1.
    else:
        logger.info(
            "check_pptx_text: include=%s, exclude=%s, extracted_text=%s",
            include, exclude, text
        )
        return 0.
