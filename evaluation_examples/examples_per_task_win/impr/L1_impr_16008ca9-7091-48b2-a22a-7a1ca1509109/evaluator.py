"""
Evaluator for task: Change the title font size on ALL slides to 36pt.

Checks the agent's output PPTX file to verify that all slide title shapes
have font size = 36pt. Does not compare against a golden file — the check
is self-contained against the task requirement.
"""
import logging

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

logger = logging.getLogger("desktopenv.metric.title_font_size")


def _find_title_shapes(slide, prs):
    """Find title shapes on a slide.

    Strategy:
    1. Look for placeholder shapes with type TITLE or CENTER_TITLE.
    2. If none found, use position heuristic: text shape near the top
       (upper third of slide height) that has text content.
    """
    title_shapes = []

    # Strategy 1: placeholder-based
    for shape in slide.shapes:
        if hasattr(shape, "is_placeholder") and shape.is_placeholder:
            try:
                ph_type = shape.placeholder_format.type
                if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    if hasattr(shape, "text_frame") and shape.has_text_frame:
                        title_shapes.append(shape)
            except Exception:
                pass

    # Strategy 2: fallback — shape near top with text
    if not title_shapes:
        top_threshold = prs.slide_height // 3
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                if shape.text and shape.text.strip():
                    if shape.top < top_threshold:
                        title_shapes.append(shape)

    return title_shapes


def _get_font_size_from_placeholder(shape):
    """Try to resolve inherited font size from the slide layout master placeholder."""
    try:
        if not (hasattr(shape, "is_placeholder") and shape.is_placeholder):
            return None
        # Access the placeholder definition in the slide layout
        slide_layout = shape.part.slide_layout
        for ph in slide_layout.placeholders:
            if ph.placeholder_format.idx == shape.placeholder_format.idx:
                if ph.has_text_frame:
                    for para in ph.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size is not None:
                                return run.font.size
                break
    except Exception:
        pass
    return None


def check_title_font_size(result_path, expected=None, **options):
    """Check that all slide titles have the specified font size.

    Args:
        result_path: Path to the agent's output PPTX file.
        expected: Not used (self-contained check). Kept for framework compatibility.
        **options:
            target_size (float): Target font size in points. Default 36.
            size_tolerance (float): Tolerance in points. Default 0.5.

    Returns:
        float: 1.0 if all title font sizes match the target, 0.0 otherwise.
    """
    target_size_pt = options.get("target_size", 36)
    size_tolerance_pt = options.get("size_tolerance", 0.5)
    target_size_emu = target_size_pt * 12700
    size_tolerance_emu = size_tolerance_pt * 12700

    try:
        prs = Presentation(result_path)
    except Exception as e:
        logger.error(f"Failed to open PPTX file '{result_path}': {e}")
        return 0.0

    if len(prs.slides) == 0:
        logger.warning("Presentation has no slides")
        return 0.0

    title_found_anywhere = False

    for slide_idx, slide in enumerate(prs.slides):
        title_shapes = _find_title_shapes(slide, prs)

        if not title_shapes:
            # Slide has no identifiable title — skip (e.g. image-only slides)
            logger.debug(f"Slide {slide_idx + 1}: no title shape found, skipping")
            continue

        title_found_anywhere = True

        for shape in title_shapes:
            shape_text = shape.text.strip()[:50] if shape.text else ""
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    actual_size = run.font.size
                    if actual_size is None:
                        # Font size inherited from layout/master — try to resolve
                        actual_size = _get_font_size_from_placeholder(shape)
                        if actual_size is None:
                            logger.debug(
                                f"Slide {slide_idx + 1}, title '{shape_text}': "
                                f"font size is None (inherited, unresolvable), skipping run"
                            )
                            continue

                    if abs(actual_size - target_size_emu) > size_tolerance_emu:
                        actual_pt = actual_size / 12700.0
                        logger.info(
                            f"Slide {slide_idx + 1}, title '{shape_text}': "
                            f"font size {actual_pt:.1f}pt != {target_size_pt}pt "
                            f"(diff {abs(actual_pt - target_size_pt):.1f}pt > "
                            f"tolerance {size_tolerance_pt}pt)"
                        )
                        return 0.0

        logger.debug(
            f"Slide {slide_idx + 1}: title font size check passed "
            f"({target_size_pt}pt)"
        )

    if not title_found_anywhere:
        logger.warning("No title shapes found on any slide")
        return 0.0

    logger.info(
        f"All slide titles have font size {target_size_pt}pt "
        f"(tolerance ±{size_tolerance_pt}pt)"
    )
    return 1.0
