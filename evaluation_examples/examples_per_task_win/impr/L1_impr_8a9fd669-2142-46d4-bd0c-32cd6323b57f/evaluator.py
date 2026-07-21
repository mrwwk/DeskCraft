import logging

from pptx import Presentation

logger = logging.getLogger("desktopenv.metric.slides")


def check_slide1_colors(result_path, expected=None, **options):
    """
    Check that slide 1 of the result PPTX has:
    - Background color = dark blue (#00008B)
    - All text on slide 1 = white (#FFFFFF)

    This evaluator only inspects slide 1 and only checks background color
    and text color, ignoring all other properties (shape dimensions, fonts,
    alignment, other slides, etc.) as they are irrelevant to this task.
    """
    prs = Presentation(result_path)

    if len(prs.slides) == 0:
        logger.error("Presentation has no slides")
        return 0

    slide = prs.slides[0]

    # --- 1. Check slide 1 background color is #00008B ---
    fill = slide.background.fill
    bg_rgb = None

    if fill.type == 1:  # solid fill on the slide itself
        try:
            bg_rgb = str(fill.fore_color.rgb)
        except Exception:
            pass
    elif fill.type == 5:  # inherited from slide master
        try:
            master = slide.slide_layout.slide_master
            master_fill = master.background.fill
            if master_fill.type == 1:
                bg_rgb = str(master_fill.fore_color.rgb)
        except Exception:
            pass

    if bg_rgb != "00008B":
        logger.info(
            "Slide 1 background color mismatch: expected 00008B, got %s", bg_rgb
        )
        return 0

    # --- 2. Check all text on slide 1 is white (#FFFFFF) ---
    def check_text_in_shape(shape):
        """Recursively check text color in shape and its sub-shapes (for GROUPs)."""
        if hasattr(shape, "text_frame"):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text and run.text.strip():
                        try:
                            if (
                                hasattr(run.font.color, "rgb")
                                and run.font.color.rgb is not None
                            ):
                                color = str(run.font.color.rgb)
                                if color != "FFFFFF":
                                    logger.info(
                                        "Text color mismatch on slide 1: "
                                        "'%s' has color %s, expected FFFFFF",
                                        run.text.strip()[:50],
                                        color,
                                    )
                                    return False
                        except Exception:
                            pass

        # Recursively descend into GROUPs and other container shapes
        if hasattr(shape, "shapes"):
            for sub_shape in shape.shapes:
                if not check_text_in_shape(sub_shape):
                    return False

        return True

    for shape in slide.shapes:
        if not check_text_in_shape(shape):
            return 0

    return 1
