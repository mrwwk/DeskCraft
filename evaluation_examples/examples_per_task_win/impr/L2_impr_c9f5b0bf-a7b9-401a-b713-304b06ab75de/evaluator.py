"""
Targeted evaluator for task: (1) slide 2 body text center-align + italic,
(2) all slides background #003366.

This replaces the original all-encompassing compare_pptx_files which checked
every shape on every slide and produced false negatives on irrelevant
decorative elements.
"""
from pptx import Presentation
from pptx.enum.text import PP_ALIGN


def compare_pptx_files(result_path, expected_path, **options):
    """
    Compare only the properties relevant to this task:
      1. All slides must have solid background color = expected_background_color.
      2. On target_slide_idx, text paragraph alignment and run italic must
         match the expected (gold) file.

    Options (passed via evaluator.options in task.json):
      - target_slide_idx (int, default 2): 1-based slide index to check
        alignment and italic against expected.
      - expected_background_color (str, default "003366"): hex RGB string
        (without #) that every slide background must match.
    """
    target_slide_idx = options.get("target_slide_idx", 2)
    expected_bg_color = options.get("expected_background_color", "003366")

    prs_result = Presentation(result_path)
    prs_expected = Presentation(expected_path)

    # -------- 1. background colour of every slide --------
    for slide in prs_result.slides:
        fill = slide.background.fill
        color_rgb = None
        if fill.type == 1:          # solid fill on the slide itself
            color_rgb = str(fill.fore_color.rgb)
        elif fill.type == 5:        # inherited from slide master
            master_fill = slide.slide_layout.slide_master.background.fill
            if master_fill.type == 1:
                color_rgb = str(master_fill.fore_color.rgb)

        if color_rgb != expected_bg_color:
            return 0.0

    # -------- 2. slide 2 alignment & italic vs expected --------
    if target_slide_idx > len(prs_result.slides) or target_slide_idx > len(prs_expected.slides):
        return 0.0

    slide_r = prs_result.slides[target_slide_idx - 1]
    slide_e = prs_expected.slides[target_slide_idx - 1]

    # only consider shapes that actually contain a text frame
    text_shapes_r = [s for s in slide_r.shapes if s.has_text_frame]
    text_shapes_e = [s for s in slide_e.shapes if s.has_text_frame]

    if len(text_shapes_r) != len(text_shapes_e):
        return 0.0

    for sr, se in zip(text_shapes_r, text_shapes_e):
        pars_r = list(sr.text_frame.paragraphs)
        pars_e = list(se.text_frame.paragraphs)

        if len(pars_r) != len(pars_e):
            return 0.0

        for pr, pe in zip(pars_r, pars_e):
            # skip empty paragraphs on both sides
            if (not pr.text or not pr.text.strip()) and (not pe.text or not pe.text.strip()):
                continue

            # --- alignment ---
            ar = pr.alignment
            ae = pe.alignment
            if ar is None:
                ar = PP_ALIGN.LEFT
            if ae is None:
                ae = PP_ALIGN.LEFT
            if ar != ae:
                return 0.0

            # --- italic on every non-empty run ---
            runs_r = [r for r in pr.runs if r.text and r.text.strip()]
            runs_e = [r for r in pe.runs if r.text and r.text.strip()]

            if len(runs_r) != len(runs_e):
                return 0.0

            for rr, re in zip(runs_r, runs_e):
                it_r = rr.font.italic
                it_e = re.font.italic
                # treat None as False (not italic)
                if (it_r or False) != (it_e or False):
                    return 0.0

    return 1.0
