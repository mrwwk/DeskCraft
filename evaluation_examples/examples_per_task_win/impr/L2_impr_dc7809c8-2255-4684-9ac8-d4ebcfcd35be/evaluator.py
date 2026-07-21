# Evaluator for task: right-align titles + add note on slide 2
# Based on slides.py compare_pptx_files, with fixes for:
# 1. LibreOffice note XML parsing (sldNum placeholder fallback)
# 2. Target-based alignment check (expected file may not have right-aligned titles)
# 3. Optional text_shapes count guard
import logging
import xml.etree.ElementTree as ET
import zipfile

from pptx import Presentation

logger = logging.getLogger("desktopenv.metric.slides")
debug_logger = logging.getLogger("desktopenv.metric.slides.debug")


def enable_debug_logging():
    """Enable debug logging for PPTX comparison"""
    debug_logger.setLevel(logging.DEBUG)
    if not debug_logger.handlers:
        handler = logging.StreamHandler()
        handler.setLevel(logging.DEBUG)
        formatter = logging.Formatter('[PPTX_DEBUG] %(message)s')
        handler.setFormatter(formatter)
        debug_logger.addHandler(handler)


def get_all_text_shapes(slide):
    """Recursively get all shapes that have text, including inside GROUPs."""

    def extract_text_shapes(shape):
        results = []
        if hasattr(shape, "text") and hasattr(shape, "text_frame"):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub_shape in shape.shapes:
                results.extend(extract_text_shapes(sub_shape))
        return results

    all_text_shapes = []
    for shape in slide.shapes:
        all_text_shapes.extend(extract_text_shapes(shape))
    return all_text_shapes


def get_slide_notes_xml_fallback(slide):
    """Extract notes text directly from notesSlide XML as fallback for python-pptx.

    LibreOffice Impress may store notes text in <ph type='sldNum'> placeholders
    instead of the expected <ph type='body'>, causing python-pptx's
    notes_text_frame.text to return None or raise AttributeError.
    This function parses the raw XML to extract all <a:t> text elements.
    """
    try:
        notes_slide = slide.notes_slide
        if notes_slide is None:
            return None
        # Access the underlying XML element
        notes_xml = notes_slide._element
        # Extract all <a:t> text content
        nsmap = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        }
        texts = []
        for t_elem in notes_xml.findall('.//a:t', nsmap):
            if t_elem.text:
                texts.append(t_elem.text)
        result = ''.join(texts)
        return result if result else None
    except Exception:
        return None


def get_slide_notes(slide):
    """Get slide notes text, with XML fallback for LibreOffice format."""
    try:
        notes_slide = slide.notes_slide
        if notes_slide is None:
            return None
        # Try python-pptx's built-in method first
        try:
            text = notes_slide.notes_text_frame.text
            if text is not None and text.strip():
                return text
        except (AttributeError, TypeError):
            pass
        # Fallback: parse XML directly
        return get_slide_notes_xml_fallback(slide)
    except Exception:
        return None


def compare_pptx_files(file1_path, file2_path, **options):
    """Compare two PPTX files, or check result file against alignment/note targets.

    Supports two modes:
    1. Comparison mode (default): compare file1 (result) vs file2 (expected).
    2. Target check mode: when alignment_target is set, check that all paragraphs
       in file1 have the specified alignment, and optionally check slide 2 notes.

    Options (comparison mode):
        examine_* : boolean flags for various check dimensions (default True).
        examine_text_shapes_count : guard for text shapes count check (default True).
        enable_debug : enable debug logging (default False).

    Options (target check mode, when alignment_target is set):
        alignment_target : "LEFT", "RIGHT", "CENTER", or "JUSTIFY".
        examine_note : whether to check slide 2 notes (default True).
        expected_note_text : the expected note text for slide 2.
        enable_debug : enable debug logging (default False).
    """
    from pptx.enum.text import PP_ALIGN

    prs1 = Presentation(file1_path)
    prs2 = Presentation(file2_path)

    enable_debug = options.get("enable_debug", False)
    if enable_debug:
        enable_debug_logging()
        debug_logger.debug("=== COMPARING PPTX FILES ===")
        debug_logger.debug("File 1: %s", file1_path)
        debug_logger.debug("File 2: %s", file2_path)
        debug_logger.debug("File 1 slides: %d", len(prs1.slides))
        debug_logger.debug("File 2 slides: %d", len(prs2.slides))

    # ── Target check mode ──────────────────────────────────────────────
    alignment_target_str = options.get("alignment_target", None)

    if alignment_target_str is not None:
        from pptx.enum.shapes import PP_PLACEHOLDER

        target_map = {
            "LEFT": PP_ALIGN.LEFT,
            "RIGHT": PP_ALIGN.RIGHT,
            "CENTER": PP_ALIGN.CENTER,
            "JUSTIFY": PP_ALIGN.JUSTIFY,
        }
        target_alignment = target_map.get(alignment_target_str.upper())
        if target_alignment is None:
            logger.error("Invalid alignment_target: %s", alignment_target_str)
            return 0.0

        # Collect title placeholder types from slide masters
        title_ph_types = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}

        # Check title alignment on each slide
        for slide_idx, slide in enumerate(prs1.slides):
            # Find shapes that are title placeholders
            title_shapes = []
            for shape in slide.shapes:
                if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    if ph_type in title_ph_types:
                        title_shapes.append(shape)

            # If no title placeholder found, fall back to all text shapes on slide
            if not title_shapes:
                title_shapes = get_all_text_shapes(slide)

            for shape in title_shapes:
                for para in shape.text_frame.paragraphs:
                    para_text = (para.text or "").strip()
                    if not para_text:
                        continue
                    align = para.alignment
                    if align is None:
                        align = PP_ALIGN.LEFT
                    if align != target_alignment:
                        if enable_debug:
                            align_name = getattr(align, 'name', str(align))
                            target_name = getattr(target_alignment, 'name', str(target_alignment))
                            debug_logger.debug(
                                "MISMATCH: Slide %d, text '%s', alignment %s != target %s",
                                slide_idx + 1, para_text[:50], align_name, target_name)
                        return 0.0

        # Check notes on slide 2 if requested
        examine_note = options.get("examine_note", True)
        if examine_note:
            expected_note = options.get("expected_note_text", "")
            if len(prs1.slides) < 2:
                if enable_debug:
                    debug_logger.debug("MISMATCH: Result file has fewer than 2 slides")
                return 0.0
            slide2 = prs1.slides[1]
            actual_note = get_slide_notes(slide2)
            if actual_note is None:
                if enable_debug:
                    debug_logger.debug("MISMATCH: Slide 2 has no notes (None)")
                return 0.0
            actual_note_clean = actual_note.strip()
            expected_note_clean = expected_note.strip()
            if actual_note_clean != expected_note_clean:
                if enable_debug:
                    debug_logger.debug("MISMATCH: Slide 2 notes differ")
                    debug_logger.debug("  Expected: '%s'", expected_note_clean[:100])
                    debug_logger.debug("  Actual:   '%s'", actual_note_clean[:100])
                return 0.0

        if enable_debug:
            debug_logger.debug("=== TARGET CHECK SUCCESSFUL ===")
        return 1.0

    # ── Comparison mode (original logic with fixes) ────────────────────
    def is_approximately_equal(val1, val2, tolerance=0.005):
        if val1 == val2:
            return True
        if val1 == 0 and val2 == 0:
            return True
        if val1 == 0 or val2 == 0:
            return False
        return abs(val1 - val2) / max(abs(val1), abs(val2)) <= tolerance

    def nonempty_runs(para):
        return [r for r in para.runs if (r.text or "").strip() != ""]

    examine_number_of_slides = options.get("examine_number_of_slides", True)
    examine_shape = options.get("examine_shape", True)
    examine_text = options.get("examine_text", True)
    examine_indent = options.get("examine_indent", True)
    examine_font_name = options.get("examine_font_name", True)
    examine_font_size = options.get("examine_font_size", True)
    examine_font_bold = options.get("examine_font_bold", True)
    examine_font_italic = options.get("examine_font_italic", True)
    examine_color_rgb = options.get("examine_color_rgb", True)
    examine_font_underline = options.get("examine_font_underline", True)
    examine_strike_through = options.get("examine_strike_through", True)
    examine_alignment = options.get("examine_alignment", True)
    examine_title_bottom_position = options.get("examine_title_bottom_position", False)
    examine_table_bottom_position = options.get("examine_table_bottom_position", False)
    examine_right_position = options.get("examine_right_position", False)
    examine_top_position = options.get("examine_top_position", False)
    examine_shape_for_shift_size = options.get("examine_shape_for_shift_size", False)
    examine_image_size = options.get("examine_image_size", False)
    examine_modify_height = options.get("examine_modify_height", False)
    examine_bullets = options.get("examine_bullets", True)
    examine_background_color = options.get("examine_background_color", True)
    examine_note = options.get("examine_note", True)
    examine_text_shapes_count = options.get("examine_text_shapes_count", True)

    if len(prs1.slides) != len(prs2.slides) and examine_number_of_slides:
        if enable_debug:
            debug_logger.debug("MISMATCH: Number of slides differ - File1: %d, File2: %d",
                               len(prs1.slides), len(prs2.slides))
        return 0.0

    slide_idx = 0
    for slide1, slide2 in zip(prs1.slides, prs2.slides):
        slide_idx += 1
        if enable_debug:
            debug_logger.debug("--- Comparing Slide %d ---", slide_idx)
            debug_logger.debug("Slide %d - Shapes count: File1=%d, File2=%d",
                               slide_idx, len(slide1.shapes), len(slide2.shapes))

        # Notes check (fixed get_slide_notes)
        if examine_note:
            note1 = get_slide_notes(slide1)
            note2 = get_slide_notes(slide2)
            note1_clean = (note1 or "").strip()
            note2_clean = (note2 or "").strip()
            if note1_clean != note2_clean:
                if enable_debug:
                    debug_logger.debug("    MISMATCH: Slide %d - Notes differ:", slide_idx)
                    debug_logger.debug("      Notes1: '%s'", note1_clean[:100])
                    debug_logger.debug("      Notes2: '%s'", note2_clean[:100])
                return 0.0

        text_shapes1 = get_all_text_shapes(slide1)
        text_shapes2 = get_all_text_shapes(slide2)

        if enable_debug:
            debug_logger.debug("Slide %d - Text shapes found: File1=%d, File2=%d",
                               slide_idx, len(text_shapes1), len(text_shapes2))

        if len(slide1.shapes) != len(slide2.shapes):
            if enable_debug:
                debug_logger.debug("MISMATCH: Slide %d - Different number of shapes: File1=%d, File2=%d",
                                   slide_idx, len(slide1.shapes), len(slide2.shapes))
            return 0.0

        shape_idx = 0
        for shape1, shape2 in zip(slide1.shapes, slide2.shapes):
            shape_idx += 1
            if enable_debug:
                debug_logger.debug("  Shape %d - Type: %s vs %s", shape_idx, shape1.shape_type, shape2.shape_type)
                if hasattr(shape1, "text") and hasattr(shape2, "text"):
                    debug_logger.debug("  Shape %d - Text: '%s' vs '%s'", shape_idx,
                                       shape1.text.strip()[:50], shape2.text.strip()[:50])
                    debug_logger.debug("  Shape %d - Position: (%d, %d) vs (%d, %d)",
                                       shape_idx, shape1.left, shape1.top, shape2.left, shape2.top)
                    debug_logger.debug("  Shape %d - Size: (%d, %d) vs (%d, %d)",
                                       shape_idx, shape1.width, shape1.height, shape2.width, shape2.height)

            if examine_title_bottom_position:
                if hasattr(shape1, "text") and hasattr(shape2, "text") and shape1.text == shape2.text:
                    if shape1.text == "Product Comparison" and (shape1.top <= shape2.top or shape1.top < 3600000):
                        return 0.0
                elif (not is_approximately_equal(shape1.left, shape2.left)
                      or not is_approximately_equal(shape1.top, shape2.top)
                      or not is_approximately_equal(shape1.width, shape2.width)
                      or not is_approximately_equal(shape1.height, shape2.height)):
                    return 0.0

            if examine_table_bottom_position:
                if slide_idx == 3 and shape1.shape_type == 19 and shape2.shape_type == 19:
                    if shape1.top <= shape2.top or shape1.top < 3600000:
                        return 0.0
                elif (not is_approximately_equal(shape1.left, shape2.left)
                      or not is_approximately_equal(shape1.top, shape2.top)
                      or not is_approximately_equal(shape1.width, shape2.width)
                      or not is_approximately_equal(shape1.height, shape2.height)):
                    return 0.0

            if examine_right_position:
                if slide_idx == 2 and not hasattr(shape1, "text") and not hasattr(shape2, "text"):
                    if shape1.left <= shape2.left or shape1.left < 4320000:
                        return 0.0

            if examine_top_position:
                if slide_idx == 2 and shape1.shape_type == 13 and shape2.shape_type == 13:
                    if shape1.top >= shape2.top or shape1.top > 1980000:
                        return 0.0

            if examine_shape_for_shift_size:
                if (not is_approximately_equal(shape1.left, shape2.left)
                    or not is_approximately_equal(shape1.top, shape2.top)
                    or not is_approximately_equal(shape1.width, shape2.width)
                    or not is_approximately_equal(shape1.height, shape2.height)):
                    if not (hasattr(shape1, "text") and hasattr(shape2, "text")
                            and shape1.text == shape2.text
                            and shape1.text == "Elaborate on what you want to discuss."):
                        return 0.0

            if ((not is_approximately_equal(shape1.left, shape2.left)
                 or not is_approximately_equal(shape1.top, shape2.top)
                 or not is_approximately_equal(shape1.width, shape2.width)
                 or not is_approximately_equal(shape1.height, shape2.height)) and examine_shape):
                if enable_debug:
                    debug_logger.debug("    MISMATCH: Slide %d, Shape %d - Shape dimensions differ:",
                                       slide_idx, shape_idx)
                    debug_logger.debug("      Left: %d vs %d", shape1.left, shape2.left)
                    debug_logger.debug("      Top: %d vs %d", shape1.top, shape2.top)
                    debug_logger.debug("      Width: %d vs %d", shape1.width, shape2.width)
                    debug_logger.debug("      Height: %d vs %d", shape1.height, shape2.height)
                return 0.0

            if examine_image_size:
                if shape1.shape_type == 13 and shape2.shape_type == 13:
                    if not is_approximately_equal(shape1.width, shape2.width) or not is_approximately_equal(shape1.height, shape2.height):
                        return 0.0
                elif (not is_approximately_equal(shape1.left, shape2.left)
                      or not is_approximately_equal(shape1.top, shape2.top)
                      or not is_approximately_equal(shape1.width, shape2.width)
                      or not is_approximately_equal(shape1.height, shape2.height)):
                    return 0.0

            if examine_modify_height:
                if not hasattr(shape1, "text") and not hasattr(shape2, "text") or shape1.shape_type == 5 and shape2.shape_type == 5:
                    if not is_approximately_equal(shape1.height, shape2.height):
                        return 0.0
                elif (not is_approximately_equal(shape1.left, shape2.left)
                      or not is_approximately_equal(shape1.top, shape2.top)
                      or not is_approximately_equal(shape1.width, shape2.width)
                      or not is_approximately_equal(shape1.height, shape2.height)):
                    return 0.0

            if shape1.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                table1 = shape1.table
                table2 = shape2.table
                if len(table1.rows) != len(table2.rows) or len(table1.columns) != len(table2.columns):
                    return 0.0

            if hasattr(shape1, "text") and hasattr(shape2, "text"):
                if shape1.text.strip() != shape2.text.strip() and examine_text:
                    return 0.0

                if len(shape1.text_frame.paragraphs) != len(shape2.text_frame.paragraphs):
                    if enable_debug:
                        debug_logger.debug("    MISMATCH: Slide %d, Shape %d - Different number of paragraphs: %d vs %d",
                                           slide_idx, shape_idx,
                                           len(shape1.text_frame.paragraphs),
                                           len(shape2.text_frame.paragraphs))
                    return 0.0

                para_idx = 0
                for para1, para2 in zip(shape1.text_frame.paragraphs, shape2.text_frame.paragraphs):
                    para_idx += 1
                    if examine_alignment:
                        align1 = para1.alignment
                        align2 = para2.alignment

                        if enable_debug:
                            align1_name = "None" if align1 is None else getattr(align1, 'name', str(align1))
                            align2_name = "None" if align2 is None else getattr(align2, 'name', str(align2))
                            debug_logger.debug("    Slide %d, Shape %d, Para %d - Alignment: '%s' vs '%s'",
                                               slide_idx, shape_idx, para_idx, align1_name, align2_name)
                            debug_logger.debug("    Slide %d, Shape %d, Para %d - Text: '%s' vs '%s'",
                                               slide_idx, shape_idx, para_idx, para1.text, para2.text)

                        if align1 is None:
                            align1 = PP_ALIGN.LEFT
                        if align2 is None:
                            align2 = PP_ALIGN.LEFT

                        if align1 != align2:
                            if enable_debug:
                                align1_final = getattr(align1, 'name', str(align1))
                                align2_final = getattr(align2, 'name', str(align2))
                                debug_logger.debug("    MISMATCH: Slide %d, Shape %d, Para %d - Alignment differs: '%s' vs '%s'",
                                                   slide_idx, shape_idx, para_idx, align1_final, align2_final)
                            return 0.0

                    if para1.text != para2.text and examine_text:
                        return 0.0

                    if para1.level != para2.level and examine_indent:
                        return 0.0

                    runs1 = para1.runs
                    runs2 = para2.runs
                    if (para1.text or "").strip() == "" and (para2.text or "").strip() == "":
                        runs1 = nonempty_runs(para1)
                        runs2 = nonempty_runs(para2)

                    if len(runs1) != len(runs2):
                        if enable_debug:
                            debug_logger.debug(
                                "    MISMATCH: Slide %d, Shape %d, Para %d - Different number of runs: %d vs %d",
                                slide_idx, shape_idx, para_idx, len(runs1), len(runs2))
                        return 0.0

                    for run1, run2 in zip(runs1, runs2):
                        if run1.font.name != run2.font.name and examine_font_name:
                            return 0.0
                        if run1.font.size != run2.font.size and examine_font_size:
                            return 0.0
                        if run1.font.bold != run2.font.bold and examine_font_bold:
                            if not ((run1.font.bold is None or run1.font.bold is False)
                                    and (run2.font.bold is None or run2.font.bold is False)):
                                return 0.0
                        if run1.font.italic != run2.font.italic and examine_font_italic:
                            if not ((run1.font.italic is None or run1.font.italic is False)
                                    and (run2.font.italic is None or run2.font.italic is False)):
                                return 0.0
                        if hasattr(run1.font.color, "rgb") and hasattr(run2.font.color, "rgb"):
                            if run1.font.color.rgb != run2.font.color.rgb and examine_color_rgb:
                                return 0.0
                        if run1.font.underline != run2.font.underline and examine_font_underline:
                            if run1.font.underline is not None and run2.font.underline is not None:
                                return 0.0
                            if ((run1.font.underline is None and run2.font.underline is True)
                                    or (run1.font.underline is True and run2.font.underline is None)):
                                return 0.0
                        if run1.font._element.attrib.get('strike', 'noStrike') != run2.font._element.attrib.get(
                                'strike', 'noStrike') and examine_strike_through:
                            return 0.0

                        # Bullet check (simplified - skip on error)
                        if examine_bullets:
                            try:
                                def _extract_bullets(xml_data):
                                    root = ET.fromstring(xml_data)
                                    ns = {
                                        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                                        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                                    }
                                    bullets = []
                                    for paragraph in root.findall('.//a:p', ns):
                                        pPr = paragraph.find('a:pPr', ns)
                                        if pPr is not None:
                                            lvl = pPr.get('lvl')
                                            buChar = pPr.find('a:buChar', ns)
                                            char = buChar.get('char') if buChar is not None else "No Bullet"
                                            buClr = pPr.find('a:buClr/a:srgbClr', ns)
                                            color = buClr.get('val') if buClr is not None else "No Color"
                                        else:
                                            lvl = "No Level"
                                            char = "No Bullet"
                                            color = "No Color"
                                        text = "".join(t.text for t in paragraph.findall('.//a:t', ns))
                                        if text.strip():
                                            bullets.append((lvl, char, text, color))
                                    return bullets

                                bullets1 = _extract_bullets(run1.part.blob.decode('utf-8'))
                                bullets2 = _extract_bullets(run2.part.blob.decode('utf-8'))
                                if len(bullets1) != len(bullets2):
                                    return 0.0
                                for (lvl1, char1, text1, _), (lvl2, char2, text2, _) in zip(bullets1, bullets2):
                                    if text1 != text2 or char1 != char2:
                                        return 0.0
                                    n_lvl1 = '0' if lvl1 is None else lvl1
                                    n_lvl2 = '0' if lvl2 is None else lvl2
                                    if n_lvl1 != n_lvl2:
                                        return 0.0
                            except Exception:
                                pass

        # Additional text shapes check (with option guard)
        if examine_alignment and len(text_shapes1) == len(text_shapes2):
            for idx, (tshape1, tshape2) in enumerate(zip(text_shapes1, text_shapes2)):
                if tshape1.text.strip() != tshape2.text.strip() and examine_text:
                    return 0.0
                if len(tshape1.text_frame.paragraphs) != len(tshape2.text_frame.paragraphs):
                    return 0.0
                for para1, para2 in zip(tshape1.text_frame.paragraphs, tshape2.text_frame.paragraphs):
                    align1 = para1.alignment
                    align2 = para2.alignment
                    if align1 is None:
                        align1 = PP_ALIGN.LEFT
                    if align2 is None:
                        align2 = PP_ALIGN.LEFT
                    if align1 != align2:
                        return 0.0
        elif examine_text_shapes_count and len(text_shapes1) != len(text_shapes2):
            if enable_debug:
                debug_logger.debug("MISMATCH: Different number of text shapes - %d vs %d",
                                   len(text_shapes1), len(text_shapes2))
            return 0.0

    if enable_debug:
        debug_logger.debug("=== COMPARISON SUCCESSFUL - Files match ===")
    return 1.0
