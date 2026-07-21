import logging

from pptx import Presentation

logger = logging.getLogger("desktopenv.metric.slides")


def _parse_hex_color(hex_str):
    """Parse a hex color string (with or without #) to an RGB tuple (R, G, B)."""
    hex_str = hex_str.lstrip("#").upper().strip()
    if len(hex_str) != 6:
        return None
    try:
        return tuple(int(hex_str[i:i + 2], 16) for i in range(0, 6, 2))
    except ValueError:
        return None


def _get_slide_background_color_rgb(slide):
    """Get the effective background color of a slide as an RGB tuple, or None if indeterminate."""
    fill = slide.background.fill
    try:
        if fill.type == 1:  # MSO_FILL_TYPE.SOLID
            if hasattr(fill.fore_color, "rgb") and fill.fore_color.rgb is not None:
                return fill.fore_color.rgb
        elif fill.type == 5:  # MSO_FILL_TYPE.BACKGROUND (inherit from master)
            master = slide.slide_layout.slide_master
            master_fill = master.background.fill
            if master_fill.type == 1:
                if hasattr(master_fill.fore_color, "rgb") and master_fill.fore_color.rgb is not None:
                    return master_fill.fore_color.rgb
    except Exception:
        pass
    return None


def _normalize_text(text):
    return " ".join((text or "").split())


def _find_title_shape(slide, target_title_text=None):
    """Find the title shape on a slide.

    Uses the expected title text first because this deck has a top chapter
    label above the real title. Then tries slide.shapes.title.
    Falls back to the topmost text-containing shape.
    """
    normalized_target = _normalize_text(target_title_text)
    if normalized_target:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and _normalize_text(shape.text) == normalized_target:
                return shape

    # Try the official title placeholder
    try:
        title = slide.shapes.title
        if title is not None and hasattr(title, "text_frame") and title.text.strip():
            return title
    except Exception:
        pass

    # Fallback: collect all text shapes and pick the topmost one
    text_shapes = []
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text.strip():
            text_shapes.append(shape)

    if not text_shapes:
        return None

    # Title is typically at the top; sort by (top, left)
    text_shapes.sort(key=lambda s: (s.top if s.top is not None else 0,
                                     s.left if s.left is not None else 0))
    return text_shapes[0]


def compare_pptx_files(result_path, expected_path=None, **options):
    """Validate a PPTX file against task-specific formatting requirements.

    This evaluator checks ONLY the conditions explicitly required by the task:
      - Slide 3 title: bold  &  color == #006400
      - Slide 3 speaker notes contain the required text
      - Every slide background color == #D3D3D3

    It intentionally does NOT perform structural comparison (shape counts,
    positions, sizes, font names, italic/underline/strikethrough, alignment,
    indent, bullets, etc.) so that LibreOffice round-trip differences do not
    produce false negatives.
    """
    # --- Read options -------------------------------------------------------
    examine_slide3_title_bold = options.get("examine_slide3_title_bold", True)
    examine_slide3_title_color = options.get("examine_slide3_title_color", True)
    target_title_text = options.get("target_title_text", "Forests as Biodiversity Reservoirs")
    target_title_color = options.get("target_title_color", "006400")
    examine_slide3_note = options.get("examine_slide3_note", True)
    target_note_text = options.get(
        "target_note_text",
        "Forests cover approximately 31% of the Earth's land surface and "
        "host over 80% of terrestrial biodiversity."
    )
    examine_background_color = options.get("examine_background_color", True)
    target_background = options.get("target_background", "D3D3D3")

    target_title_rgb = _parse_hex_color(target_title_color)
    target_bg_rgb = _parse_hex_color(target_background)

    if target_title_rgb is None:
        logger.error("Invalid target_title_color: %s", target_title_color)
        return 0.0
    if target_bg_rgb is None:
        logger.error("Invalid target_background: %s", target_background)
        return 0.0

    # --- Open result file ---------------------------------------------------
    try:
        prs = Presentation(result_path)
    except Exception as exc:
        logger.error("Failed to open PPTX file: %s", exc)
        return 0.0

    total_slides = len(prs.slides)
    logger.info("Opened PPTX with %d slides", total_slides)

    # ========================================================================
    # Check 1 — Slide 3 title: bold + color
    # ========================================================================
    if examine_slide3_title_bold or examine_slide3_title_color:
        if total_slides < 3:
            logger.error("Presentation has only %d slides; slide 3 required", total_slides)
            return 0.0

        slide3 = prs.slides[2]
        title_shape = _find_title_shape(slide3, target_title_text)

        if title_shape is None:
            logger.error("Could not locate a title shape on slide 3")
            return 0.0

        logger.info("Slide 3 title shape text: %s", title_shape.text.strip()[:80])

        # --- Bold ---
        if examine_slide3_title_bold:
            for para in title_shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip() and not run.font.bold:
                        logger.error(
                            "Slide 3 title: run '%s' is not bold (bold=%s)",
                            run.text[:40], run.font.bold,
                        )
                        return 0.0
            logger.info("Slide 3 title bold: PASSED")

        # --- Color ---
        if examine_slide3_title_color:
            for para in title_shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    try:
                        run_rgb = run.font.color.rgb
                    except Exception as exc:
                        logger.error(
                            "Slide 3 title: cannot read color for run '%s': %s",
                            run.text[:40], exc,
                        )
                        return 0.0
                    if run_rgb != target_title_rgb:
                        logger.error(
                            "Slide 3 title: run '%s' color %s != expected %s",
                            run.text[:40], run_rgb, target_title_rgb,
                        )
                        return 0.0
            logger.info("Slide 3 title color: PASSED")

    # ========================================================================
    # Check 2 — Slide 3 speaker notes
    # ========================================================================
    if examine_slide3_note:
        if total_slides < 3:
            logger.error("Presentation has only %d slides; cannot check slide 3 notes", total_slides)
            return 0.0

        slide3 = prs.slides[2]
        try:
            notes_slide = slide3.notes_slide
            notes_text = notes_slide.notes_text_frame.text if notes_slide else ""
        except Exception:
            notes_text = ""

        logger.info("Slide 3 notes text (first 120 chars): %s", notes_text.strip()[:120])

        if target_note_text.strip() not in notes_text:
            logger.error("Slide 3 notes: required text not found")
            logger.error("  Expected: %s", target_note_text.strip())
            logger.error("  Actual:   %s", notes_text.strip())
            return 0.0

        logger.info("Slide 3 notes: PASSED")

    # ========================================================================
    # Check 3 — Background color on every slide
    # ========================================================================
    if examine_background_color:
        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            bg_rgb = _get_slide_background_color_rgb(slide)
            if bg_rgb is None:
                logger.error("Slide %d: could not determine background color", slide_num)
                return 0.0
            if bg_rgb != target_bg_rgb:
                logger.error(
                    "Slide %d: background color %s != expected %s",
                    slide_num, bg_rgb, target_bg_rgb,
                )
                return 0.0
            logger.info("Slide %d background: PASSED (%s)", slide_num, bg_rgb)

    logger.info("All task-specific checks PASSED")
    return 1.0
