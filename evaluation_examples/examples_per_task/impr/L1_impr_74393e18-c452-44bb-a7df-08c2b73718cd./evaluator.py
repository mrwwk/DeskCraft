"""
Evaluator for: "Make all title text bold AND italic across ALL slides."

Checks every slide for title shapes and verifies that all text runs in those
titles have bold=True and italic=True.  Only title placeholder shapes
(TITLE / CENTER_TITLE) are examined; non-title text, shape dimensions,
background colours, notes and other unrelated properties are intentionally
ignored so that the evaluator is robust against LibreOffice re-rendering
artifacts and master-slide side-effects.
"""
import logging

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

logger = logging.getLogger("desktopenv.metric.slides")


def check_all_titles_bold_italic(result_path, expected=None, **options):
    """Return 1.0 when every title run on every slide is bold AND italic.

    * ``result_path`` – local path to the saved PPTX (from vm_file getter).
    * ``expected``   – not used (no golden-file comparison needed).
    * ``**options``  – reserved for future tuning (currently unused).

    A slide without any title placeholder is skipped (vacuously satisfied).
    The function returns **0.0** as soon as a single title run fails the
    bold+italic check.
    """
    prs = Presentation(result_path)

    found_any_title = False

    for slide_idx, slide in enumerate(prs.slides):
        title_shapes = _find_title_shapes(slide)

        if not title_shapes:
            # No title placeholder on this slide → nothing to verify
            continue

        for shape in title_shapes:
            found_any_title = True

            if not shape.has_text_frame:
                continue

            text = shape.text_frame.text.strip()
            if not text:
                # Empty title placeholder → vacuously satisfied
                continue

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not (run.text or "").strip():
                        continue

                    if run.font.bold is not True:
                        logger.info(
                            "Slide %d title run bold check failed: "
                            "bold=%s (expected True), text='%s'",
                            slide_idx + 1, run.font.bold,
                            run.text[:80],
                        )
                        return 0.0

                    if run.font.italic is not True:
                        logger.info(
                            "Slide %d title run italic check failed: "
                            "italic=%s (expected True), text='%s'",
                            slide_idx + 1, run.font.italic,
                            run.text[:80],
                        )
                        return 0.0

    if not found_any_title:
        logger.warning("No title placeholder shapes found in the presentation")
        return 0.0

    logger.info("All title runs on all slides are bold=True AND italic=True")
    return 1.0


def _find_title_shapes(slide):
    """Return a list of title placeholder shapes on *slide*.

    If no TITLE or CENTER_TITLE placeholder is found, fall back to the
    first text-containing shape on the slide (common in LibreOffice-produced
    PPTX files where the placeholder metadata may be lost).
    """
    title_shapes = []

    for shape in slide.shapes:
        if hasattr(shape, "is_placeholder") and shape.is_placeholder:
            try:
                ph = shape.placeholder_format
                if ph.type in (PP_PLACEHOLDER.TITLE,
                               PP_PLACEHOLDER.CENTER_TITLE):
                    title_shapes.append(shape)
            except (ValueError, AttributeError):
                pass

    if title_shapes:
        return title_shapes

    # Fallback: first non-empty text shape on the slide
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.has_text_frame:
            if shape.text_frame.text.strip():
                return [shape]

    return []
