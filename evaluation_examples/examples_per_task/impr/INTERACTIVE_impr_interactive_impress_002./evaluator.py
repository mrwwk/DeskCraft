import logging
import os
from typing import Dict, List, Optional

logger = logging.getLogger("desktopenv.metric.general")


def check_include_exclude(result: str, rules: Dict[str, List[str]]) -> float:
    """
    Check if the PPTX file contains all required strings and excludes forbidden ones.

    The result parameter is a local file path to the PPTX (pulled from VM by vm_file getter).
    This evaluator reads the PPTX using python-pptx, extracts all text, and checks
    include/exclude rules against the extracted text.

    Args:
        result: Path to PPTX file (from vm_file getter) or raw text string (fallback).
        rules: Dict with "include" and "exclude" lists of strings.

    Returns:
        1.0 if all include strings are found and no exclude strings are found, else 0.0.
    """
    if result is None:
        logger.warning("Result is None, returning 0.0")
        return 0.0

    include = rules.get("include", [])
    exclude = rules.get("exclude", [])

    # Attempt to read result as a PPTX file
    text = _extract_pptx_text(result)

    # If PPTX extraction failed, fall back to treating result as plain text
    if text is None:
        text = str(result) if result else ""

    if all(r in text for r in include) and all(r not in text for r in exclude):
        return 1.0
    else:
        logger.info(
            f"check_include_exclude: include={include}, exclude={exclude}, "
            f"text_preview={text[:200] if text else '(empty)'}"
        )
        return 0.0


def _extract_pptx_text(file_path: str) -> Optional[str]:
    """
    Extract all text from a PPTX file using python-pptx.

    Returns the concatenated text of all slides, or None if extraction fails.
    """
    if not os.path.isfile(file_path):
        logger.warning(f"File not found: {file_path}")
        return None

    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text:
                                text_parts.append(run.text)
        return ' '.join(text_parts)
    except Exception as e:
        logger.warning(f"Failed to read PPTX from {file_path}: {e}")
        return None
