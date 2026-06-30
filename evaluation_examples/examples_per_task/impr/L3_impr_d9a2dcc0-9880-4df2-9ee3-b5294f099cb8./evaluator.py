"""
Targeted evaluator for Van Gogh presentation task.
Checks only the 3 specific task conditions on Slides 4-5,
avoiding false negatives from LibreOffice format artifacts on unrelated slides.
"""
from pptx import Presentation


def _find_shapes_by_keywords(slide, keywords):
    """Find all shapes on a slide whose text contains any of the given keywords (case-insensitive)."""
    matches = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text
        if any(kw.lower() in text.lower() for kw in keywords):
            matches.append(shape)
    return matches


def _all_text_runs(shape):
    """Yield all runs with non-empty text from a shape's text frame."""
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if (run.text or "").strip():
                yield run


def check_slide5_body_format(result_path, expected=None, **options):
    """
    Check Slide 5 body text is dark red (#8B0000) and italic.
    Returns 1.0 if both conditions are met, 0.0 otherwise.
    """
    try:
        prs = Presentation(result_path)
        target_slide = options.get("target_slide", 4)  # 0-indexed: 4 = Slide 5
        target_color = options.get("target_color_rgb", "8B0000")
        keywords = options.get("body_text_keywords", ["decade of transformation"])

        if len(prs.slides) <= target_slide:
            return 0.0

        slide = prs.slides[target_slide]
        shapes = _find_shapes_by_keywords(slide, keywords)

        if not shapes:
            return 0.0

        found_color_ok = False
        found_italic_ok = False

        for shape in shapes:
            for run in _all_text_runs(shape):
                # Check color
                if not found_color_ok:
                    if hasattr(run.font.color, 'rgb') and run.font.color.rgb is not None:
                        color_str = str(run.font.color.rgb)
                        if color_str.upper() == target_color.upper():
                            found_color_ok = True

                # Check italic
                if not found_italic_ok:
                    if run.font.italic is True:
                        found_italic_ok = True

                if found_color_ok and found_italic_ok:
                    return 1.0

        return 1.0 if (found_color_ok and found_italic_ok) else 0.0
    except Exception:
        return 0.0


def check_slide4_title_format(result_path, expected=None, **options):
    """
    Check Slide 4 title is bold and underlined.
    Returns 1.0 if both conditions are met, 0.0 otherwise.
    """
    try:
        prs = Presentation(result_path)
        target_slide = options.get("target_slide", 3)  # 0-indexed: 3 = Slide 4
        keywords = options.get("title_keywords", ["Zundert", "Canvas"])

        if len(prs.slides) <= target_slide:
            return 0.0

        slide = prs.slides[target_slide]
        shapes = _find_shapes_by_keywords(slide, keywords)

        if not shapes:
            return 0.0

        found_bold_ok = False
        found_underline_ok = False

        for shape in shapes:
            for run in _all_text_runs(shape):
                # Check bold
                if not found_bold_ok:
                    if run.font.bold is True:
                        found_bold_ok = True

                # Check underline
                if not found_underline_ok:
                    if run.font.underline is True:
                        found_underline_ok = True

                if found_bold_ok and found_underline_ok:
                    return 1.0

        return 1.0 if (found_bold_ok and found_underline_ok) else 0.0
    except Exception:
        return 0.0


def check_slide4_note_content(result_path, expected=None, **options):
    """
    Check Slide 4 has the specified note text.
    Returns 1.0 if the note contains the expected text, 0.0 otherwise.
    """
    try:
        prs = Presentation(result_path)
        target_slide = options.get("target_slide", 3)  # 0-indexed: 3 = Slide 4
        expected_note = options.get("expected_note",
            "Van Gogh began his artistic career at 27 and produced over 860 oil paintings in just 10 years.")

        if len(prs.slides) <= target_slide:
            return 0.0

        slide = prs.slides[target_slide]
        notes_slide = slide.notes_slide

        if notes_slide is None:
            return 0.0

        actual_note = notes_slide.notes_text_frame.text.strip()

        if expected_note.strip().lower() in actual_note.lower():
            return 1.0
        return 0.0
    except Exception:
        return 0.0
