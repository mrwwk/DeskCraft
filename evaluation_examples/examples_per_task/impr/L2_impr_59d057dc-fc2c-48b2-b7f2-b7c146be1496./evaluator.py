"""
Targeted evaluator for the task: On slide 1, change the title font to "Georgia" and set the font size to 40pt.

This replaces the generic compare_pptx_files evaluator which had two issues:
1. The expected file (Ecosystem_L1.pptx) had wrong title font (Sorts Mill Goudy, 54pt) instead of Georgia 40pt.
2. With empty options, all default checks were enabled, causing failure on unrelated shape dimension differences.
"""
import logging

from pptx import Presentation

logger = logging.getLogger("desktopenv.metric.slides")


def check_slide1_title_font(result_path, **options):
    """
    Check that the title shape on slide 1 has the expected font name and size.

    Args:
        result_path: Path to the PPTX file produced by the agent.
        **options:
            font_name: Expected font name (default: "Georgia").
            font_size_pt: Expected font size in points (default: 40).

    Returns:
        1.0 if the title font matches expectations, 0.0 otherwise.
    """
    expected_font_name = options.get("font_name", "Georgia")
    expected_font_size_pt = options.get("font_size_pt", 40)
    # 1pt = 12700 EMU in python-pptx
    expected_font_size_emu = expected_font_size_pt * 12700
    # Allow 1pt tolerance for LibreOffice rounding
    size_tolerance_emu = 12700

    try:
        prs = Presentation(result_path)
    except Exception as e:
        logger.error(f"Failed to open PPTX file: {e}")
        return 0.0

    if len(prs.slides) == 0:
        logger.error("Presentation has no slides")
        return 0.0

    slide = prs.slides[0]

    # Strategy 1: Try to find the title shape via placeholder type.
    # TITLE (1) / CENTER_TITLE = 1, SUBTITLE = 3
    title_shape = None
    for shape in slide.shapes:
        try:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type in (1, 3):
                    title_shape = shape
                    break
        except Exception:
            pass

    # Strategy 2: If placeholder-based detection failed, find the shape
    # with the largest font size — typically the title.
    if title_shape is None:
        candidates = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            max_font_size = 0
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip() and run.font.size and run.font.size > max_font_size:
                        max_font_size = run.font.size
            if max_font_size > 0:
                candidates.append((max_font_size, shape))

        if candidates:
            candidates.sort(key=lambda x: x[0], reverse=True)
            title_shape = candidates[0][1]

    if title_shape is None:
        logger.error("Could not identify title shape on slide 1")
        return 0.0

    # Verify font properties on all non-empty runs of the title shape
    for para in title_shape.text_frame.paragraphs:
        for run in para.runs:
            if not run.text.strip():
                continue

            if run.font.name != expected_font_name:
                logger.info(
                    "Title font name mismatch: expected '%s', got '%s'",
                    expected_font_name, run.font.name
                )
                return 0.0

            if run.font.size is None or abs(run.font.size - expected_font_size_emu) > size_tolerance_emu:
                actual_pt = run.font.size / 12700.0 if run.font.size else None
                logger.info(
                    "Title font size mismatch: expected %dpt (%d EMU), got %spt (%s EMU)",
                    expected_font_size_pt, expected_font_size_emu,
                    actual_pt, run.font.size
                )
                return 0.0

    return 1.0
