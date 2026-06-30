"""
Evaluator for task: Open "work_summary_template_english.pptx".
On slide 1, make the title bold, 48pt, and underlined.

This evaluator checks ONLY the three required formatting properties
(bold, 48pt font size, underline) on slide 1's title shape.
It does NOT check background color, shape positions, shape types,
slide count, notes, bullets, alignment, or any other properties
that are irrelevant to this task.
"""
import logging

from pptx import Presentation

logger = logging.getLogger("desktopenv.metric.slides")


def _identify_title_shape(slide):
    """
    Identify the title shape on a slide.

    Strategy (in priority order):
    1. Use slide.shapes.title if available (title placeholder).
    2. Fall back: find the text shape with the largest font size.
    """
    # Strategy 1: title placeholder
    try:
        if hasattr(slide.shapes, 'title'):
            title = slide.shapes.title
            if title is not None and hasattr(title, "text") and title.text.strip():
                logger.info("Title identified via slide.shapes.title placeholder")
                return title
    except Exception:
        pass

    # Strategy 2: largest font size among text shapes
    candidates = []
    for shape in slide.shapes:
        if not hasattr(shape, "text") or not hasattr(shape, "text_frame"):
            continue
        text = shape.text.strip() if shape.text else ""
        if not text:
            continue
        max_size = 0
        try:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is not None and run.font.size > max_size:
                        max_size = run.font.size
        except Exception:
            continue
        if max_size > 0:
            candidates.append((shape, max_size))

    if not candidates:
        logger.warning("No text shapes with font size found on slide")
        return None

    candidates.sort(key=lambda x: x[1], reverse=True)
    best = candidates[0][0]
    logger.info(
        "Title identified via largest font size: '%s...' (max_font=%d EMU)",
        (best.text or "")[:80],
        candidates[0][1],
    )
    return best


def check_slide1_title_format(result_path, rules, **options):
    """
    Check that slide 1's title has the required formatting.

    Args:
        result_path: Path to the agent-modified PPTX file.
        rules: dict from expected getter (type=rule) with keys:
            - slide_index (int, default 0): 0-based slide index.
            - expected_font_size (int, default 4800): font size in EMU.
            - check_bold (bool, default True): whether to verify bold.
            - check_underline (bool, default True): whether to verify underline.
        options: Additional keyword options (unused).

    Returns:
        1.0 if all formatting conditions are met on all non-empty runs
        of the identified title shape; 0.0 otherwise.
    """
    if rules is None:
        rules = {}

    slide_index = rules.get("slide_index", 0)
    expected_font_size = rules.get("expected_font_size", 4800)
    check_bold = rules.get("check_bold", True)
    check_underline = rules.get("check_underline", True)

    logger.info(
        "check_slide1_title_format: slide=%d, font_size=%d EMU, bold=%s, underline=%s",
        slide_index,
        expected_font_size,
        check_bold,
        check_underline,
    )

    # Open the presentation
    try:
        prs = Presentation(result_path)
    except Exception as e:
        logger.error("Failed to open presentation '%s': %s", result_path, e)
        return 0.0

    if slide_index >= len(prs.slides):
        logger.error(
            "Slide index %d out of range (total slides: %d)",
            slide_index,
            len(prs.slides),
        )
        return 0.0

    slide = prs.slides[slide_index]

    # Identify the title shape
    title_shape = _identify_title_shape(slide)
    if title_shape is None:
        logger.error("Could not identify title shape on slide %d", slide_index)
        return 0.0

    logger.info("Checking title shape text: '%s'", (title_shape.text or "")[:120])

    # Check all non-empty runs in the title shape
    all_ok = True
    run_count = 0

    try:
        for para in title_shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text or not run.text.strip():
                    continue
                run_count += 1

                # ---- Bold check ----
                if check_bold:
                    is_bold = run.font.bold
                    # None and False both mean "not bold"
                    if is_bold is None or is_bold is False:
                        logger.info(
                            "Run '%s' is NOT bold (bold=%s)",
                            run.text[:50],
                            is_bold,
                        )
                        all_ok = False

                # ---- Font size check (48pt = 4800 EMU) ----
                actual_size = run.font.size
                if actual_size != expected_font_size:
                    logger.info(
                        "Run '%s' has font size %s (expected %d)",
                        run.text[:50],
                        actual_size,
                        expected_font_size,
                    )
                    all_ok = False

                # ---- Underline check ----
                if check_underline:
                    is_underline = run.font.underline
                    # None and False both mean "not underlined"
                    if is_underline is None or is_underline is False:
                        logger.info(
                            "Run '%s' is NOT underlined (underline=%s)",
                            run.text[:50],
                            is_underline,
                        )
                        all_ok = False
    except Exception as e:
        logger.error("Error while checking runs: %s", e)
        return 0.0

    if run_count == 0:
        logger.error("No non-empty runs found in title shape")
        return 0.0

    if all_ok:
        logger.info(
            "PASS: All %d runs are bold, %dpt, and underlined",
            run_count,
            expected_font_size // 100,
        )
        return 1.0
    else:
        logger.info("FAIL: Some runs did not meet the formatting requirements")
        return 0.0
